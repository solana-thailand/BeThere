#!/usr/bin/env python3
"""
Diagnostic: generate Keynote isolation variants.

Generates 5 progressively richer .pptx files in .deliverables/diag_*.pptx
so the user can test which one(s) Keynote rejects. The pattern of failures
pinpoints the exact structural culprit.

Output files (all Keynote-tested by user):
  diag_1_minimal.pptx    — empty Presentation() (just python-pptx template)
  diag_2_textbox.pptx    — 1 slide w/ textbox + card + accent stripe (no image)
  diag_3_image.pptx      — 1 slide w/ an embedded PNG image
  diag_4_notes.pptx      — 1 slide w/ speaker notes attached
  diag_5_full_no_notes.pptx — full 17-slide deck WITHOUT speaker notes
  diag_6_full.pptx       — full 17-slide deck WITH notes (= current behavior)

Interpretation:
  1 fails    → python-pptx base template is incompatible w/ Keynote version
  1 ok, 2 fails  → card/stripe/textbox shape patterns break Keynote
  1 ok, 3 fails  → embedded PNG image breaks Keynote
  1 ok, 4 fails  → notes slides break Keynote
  1-4 ok, 5 fails → some slide's content breaks it (binary-search next)
  1-5 ok, 6 fails → notes slides break it (confirmed)
  all ok         → the previously shipped file was stale; you're good

Run:
    python3 scripts/diag_keynote.py
"""

import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / ".deliverables"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# Load the real generator module (without invoking build()) so we can reuse
# its slide_* functions for the full-deck variants.
_spec = importlib.util.spec_from_file_location(
    "mkpd", str(ROOT / "scripts" / "make_pitch_deck.py")
)
mk = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(mk)


def _save(prs: Presentation, name: str) -> None:
    """Save, stripping printerSettings first (already-fixed hygiene)."""
    for rel in list(prs.part.rels.values()):
        if rel.reltype.endswith("/printerSettings"):
            prs.part.drop_rel(rel.rId)
    path = OUT_DIR / name
    prs.save(path)
    print(
        f"  ✓ {name}  ({path.stat().st_size:,} bytes, {len(prs.slides._sldIdLst)} slides)"
    )


def variant_1_minimal() -> None:
    """Empty Presentation — just python-pptx's default template."""
    print("[1] minimal (template only)")
    _save(Presentation(), "diag_1_minimal.pptx")


def variant_2_textbox() -> None:
    """One blank slide with the deck's shape patterns: card + accent stripe + textbox."""
    print("[2] textbox + card + accent stripe")
    prs = Presentation()
    prs.slide_width = mk.SLIDE_W
    prs.slide_height = mk.SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    mk.set_bg(s, mk.BG_DARK)
    mk.card(
        s,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(2),
        fill=mk.BG_CARD,
        border=mk.DIVIDER,
    )
    mk.add_rect(s, Inches(1), Inches(1), Inches(0.08), Inches(2), fill=mk.GREEN)
    mk.add_text(
        s,
        Inches(1.3),
        Inches(1.3),
        Inches(3),
        Inches(0.5),
        "Hello Keynote",
        size=28,
        color=mk.TEXT_LIGHT,
        bold=True,
    )
    mk.add_multi(
        s,
        Inches(1.3),
        Inches(2.0),
        Inches(3),
        Inches(0.8),
        runs=[
            [("Bold title", 17, mk.TEXT_LIGHT, True)],
            [("muted subtitle", 13, mk.TEXT_MUTED, False)],
        ],
        anchor=mk.MSO_ANCHOR.MIDDLE,
    )
    _save(prs, "diag_2_textbox.pptx")


def variant_3_image() -> None:
    """One slide with an embedded PNG (tests image parts + relationships)."""
    print("[3] embedded PNG image")
    prs = Presentation()
    prs.slide_width = mk.SLIDE_W
    prs.slide_height = mk.SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    mk.set_bg(s, mk.BG_DARK)
    s.shapes.add_picture(str(mk.LOGO_PATH), Inches(1), Inches(1), width=Inches(4))
    s.shapes.add_picture(str(mk.QR_PATH), Inches(6), Inches(1), width=Inches(3))
    _save(prs, "diag_3_image.pptx")


def variant_4_notes() -> None:
    """One slide WITH speaker notes (tests notesSlide parts + notesMaster)."""
    print("[4] speaker notes")
    prs = Presentation()
    prs.slide_width = mk.SLIDE_W
    prs.slide_height = mk.SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    mk.set_bg(s, mk.BG_DARK)
    mk.add_text(
        s,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(1),
        "With notes",
        size=36,
        color=mk.TEXT_LIGHT,
        bold=True,
    )
    s.notes_slide.notes_text_frame.text = "These are speaker notes for slide 1."
    _save(prs, "diag_4_notes.pptx")


def _clean_notes(prs: Presentation) -> None:
    """Strip non-essential elements from notes that may trip Keynote's importer.

    Removes from the notesMaster AND every notesSlide:
    - <p:extLst> descendants (p14:creationId PowerPoint 2010 extensions)
    - placeholder shapes of type hdr / dt / ftr / sldNum (keeps sldImg + body)
    """
    from pptx.oxml.ns import qn

    drop_types = {"hdr", "dt", "ftr", "sldNum"}

    def _strip(elem) -> None:
        # Remove extLst descendants (PowerPoint 2010 creationId extensions)
        for ext_lst in list(elem.iter(qn("p:extLst"))):
            ext_lst.getparent().remove(ext_lst)
        # Remove non-essential placeholder shapes (keep sldImg + body)
        for sp in list(elem.iter(qn("p:sp"))):
            nvSpPr = sp.find(qn("p:nvSpPr"))
            if nvSpPr is None:
                continue
            nvPr = nvSpPr.find(qn("p:nvPr"))
            if nvPr is None:
                continue
            ph = nvPr.find(qn("p:ph"))
            if ph is not None and ph.get("type") in drop_types:
                sp.getparent().remove(sp)

    # Clean the notesMaster (presentation-level relationship)
    for rel in prs.part.rels.values():
        if rel.reltype.endswith("/notesMaster"):
            _strip(rel.target_part._element)

    # Clean each notesSlide
    for slide in prs.slides:
        if slide.has_notes_slide:
            _strip(slide.notes_slide._element)


def variant_4a_clean_notes() -> None:
    """One slide with notes, but notesMaster + notesSlide stripped of
    non-essential elements (extLst, hdr/dt/ftr/sldNum placeholders).

    Tests whether python-pptx's default notesMaster content is what Keynote
    chokes on. If this imports, we can apply _clean_notes() in the generator
    and keep notes enabled everywhere.
    """
    print("[4a] speaker notes with cleaned notesMaster")
    prs = Presentation()
    prs.slide_width = mk.SLIDE_W
    prs.slide_height = mk.SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    mk.set_bg(s, mk.BG_DARK)
    mk.add_text(
        s,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(1),
        "With cleaned notes",
        size=36,
        color=mk.TEXT_LIGHT,
        bold=True,
    )
    s.notes_slide.notes_text_frame.text = "These are speaker notes for slide 1."
    _clean_notes(prs)
    _save(prs, "diag_4a_clean_notes.pptx")


def _build_full_deck(with_notes: bool) -> Presentation:
    """Replicate build() but optionally skip the notes loop."""
    global_counter_reset = mk._slide_counter
    mk._slide_counter = 0
    prs = Presentation()
    prs.slide_width = mk.SLIDE_W
    prs.slide_height = mk.SLIDE_H

    mk.slide_01_title(prs)
    mk.slide_join_live_demo(prs)
    mk.slide_02_problem(prs)
    mk.slide_03_solution(prs)
    mk.slide_04_demo_flow(prs)
    mk.slide_section_break(
        prs,
        eyebrow="Part Two",
        title="The Build",
        subtitle="Architecture · escrow · live dashboard · security",
    )
    mk.slide_05_architecture(prs)
    mk.slide_06_escrow(prs)
    mk.slide_07_dashboard(prs)
    mk.slide_08_performance(prs)
    mk.slide_09_security(prs)
    mk.slide_10_competitive(prs)
    mk.slide_11_whats_built(prs)
    mk.slide_section_break(
        prs,
        eyebrow="Part Three",
        title="The Road Ahead",
        subtitle="Roadmap · evidence · what's next",
    )
    mk.slide_12_roadmap(prs)
    mk.slide_13_evidence(prs)
    mk.slide_14_qa(prs)

    if with_notes:
        for slide, note in zip(prs.slides, mk.SLIDE_NOTES):
            slide.notes_slide.notes_text_frame.text = note

    mk._slide_counter = global_counter_reset
    return prs


def variant_5_full_no_notes() -> None:
    """Full 17-slide deck WITHOUT speaker notes."""
    print("[5] full deck, no notes")
    _save(_build_full_deck(with_notes=False), "diag_5_full_no_notes.pptx")


def variant_6_full_with_notes() -> None:
    """Full 17-slide deck WITH speaker notes (= current behavior)."""
    print("[6] full deck, with notes")
    _save(_build_full_deck(with_notes=True), "diag_6_full.pptx")


def variant_6a_full_clean_notes() -> None:
    """Full 17-slide deck WITH speaker notes, then _clean_notes() applied.

    If diag_4a_clean_notes.pptx imports AND this imports, the cleaning is
    sufficient for a real multi-slide deck → we can ship notes by default.
    If 4a imports but THIS fails, a specific slide's notes content still
    trips Keynote and needs more investigation.
    """
    print("[6a] full deck, with cleaned notes")
    prs = _build_full_deck(with_notes=True)
    _clean_notes(prs)
    _save(prs, "diag_6a_full_clean_notes.pptx")


def main() -> None:
    print(f"Generating Keynote isolation variants into {OUT_DIR}/\n")
    variant_1_minimal()
    variant_2_textbox()
    variant_3_image()
    variant_4_notes()
    variant_4a_clean_notes()
    variant_5_full_no_notes()
    variant_6_full_with_notes()
    variant_6a_full_clean_notes()
    print("\nDone. Open each in Keynote and report which ones fail.")
    print("Command:  open -a Keynote .deliverables/diag_1_minimal.pptx")
    print("Then use arrow-up in the shell to swap the filename for the rest.")


if __name__ == "__main__":
    main()
