#!/usr/bin/env python3
"""
BeThere — Shape Inventory Diagnostic
Dumps every shape per slide (type, position, size, text) so layout issues
like overlaps, off-canvas placement, or missing elements are easy to spot
without rendering the deck.

Run:
    python3 scripts/diag_shapes.py
    python3 scripts/diag_shapes.py --slide 3        # one slide (1-based)
    python3 scripts/diag_shapes.py --overlap        # flag overlaps
Input:
    .deliverables/bethere-pitch.pptx
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

DECK = Path(__file__).resolve().parent.parent / ".deliverables" / "bethere-pitch.pptx"

# Slide canvas (in inches) — matches make_pitch_deck.SLIDE_W / SLIDE_H
CANVAS_W_IN = 13.333
CANVAS_H_IN = 7.5


def emu_to_in(v) -> float:
    """Convert EMU to inches, rounded to 0.01"."""
    if v is None:
        return 0.0
    return round(float(v) / 914400.0, 2)


def shape_kind(shape) -> str:
    """Human-readable shape type label."""
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            return f"AUTO_SHAPE({shape.auto_shape_type})"
        except Exception:
            return "AUTO_SHAPE"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "PICTURE"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "TEXT_BOX"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "PLACEHOLDER"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "GROUP"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "TABLE"
    return str(st)


def shape_text(shape) -> str:
    """Best-effort single-line text excerpt (first non-empty paragraph)."""
    if not shape.has_text_frame:
        return ""
    tf = shape.text_frame
    for p in tf.paragraphs:
        line = "".join(r.text for r in p.runs).strip()
        if line:
            # Truncate long lines so the dump stays scannable
            return line if len(line) <= 70 else line[:67] + "..."
    return ""


def rect_in(shape) -> tuple[float, float, float, float]:
    """Return (left, top, width, height) in inches."""
    return (
        emu_to_in(shape.left),
        emu_to_in(shape.top),
        emu_to_in(shape.width),
        emu_to_in(shape.height),
    )


def boxes_overlap(
    a: tuple[float, float, float, float],
    b: tuple[float, float, float, float],
    tol: float = 0.02,
) -> bool:
    """Axis-aligned overlap test with a tolerance in inches."""
    al, at, aw, ah = a
    bl, bt, bw, bh = b
    # Expand each box by tol/2 so marginal touches count as overlaps
    return not (
        al + aw <= bl + tol
        or bl + bw <= al + tol
        or at + ah <= bt + tol
        or bt + bh <= at + tol
    )


def dump_slide(slide, idx: int, flag_overlap: bool) -> list[tuple[int, int]]:
    """Print one slide's shape inventory. Return list of (i, j) overlap pairs."""
    shapes = list(slide.shapes)
    print(f"\n{'=' * 78}")
    print(f"SLIDE {idx:02d}  ·  {len(shapes)} shapes")
    print(f"{'=' * 78}")
    print(f"{'#':>3}  {'TYPE':<28} {'L':>6} {'T':>6} {'W':>6} {'H':>6}  TEXT")
    print(f"{'-' * 3}  {'-' * 28} {'-' * 6} {'-' * 6} {'-' * 6} {'-' * 6}  {'-' * 30}")

    rects: list[tuple[float, float, float, float]] = []
    for n, sh in enumerate(shapes):
        l, t, w, h = rect_in(sh)
        rects.append((l, t, w, h))
        kind = shape_kind(sh)
        txt = shape_text(sh)
        # Flag off-canvas shapes
        off = ""
        if (
            l < -0.01
            or t < -0.01
            or l + w > CANVAS_W_IN + 0.01
            or t + h > CANVAS_H_IN + 0.01
        ):
            off = " ⚠OFF-CANVAS"
        print(f"{n:>3}  {kind:<28} {l:>6.2f} {t:>6.2f} {w:>6.2f} {h:>6.2f}  {txt}{off}")

    overlaps: list[tuple[int, int]] = []
    if flag_overlap:
        for i in range(len(rects)):
            for j in range(i + 1, len(rects)):
                if boxes_overlap(rects[i], rects[j]):
                    overlaps.append((i, j))
        if overlaps:
            print(f"\n  ⚠ {len(overlaps)} overlapping pair(s):")
            for i, j in overlaps:
                ki = shape_kind(shapes[i])[:20]
                kj = shape_kind(shapes[j])[:20]
                ti = shape_text(shapes[i])[:30]
                tj = shape_text(shapes[j])[:30]
                print(f"    [{i}] {ki} '{ti}'  ⟷  [{j}] {kj} '{tj}'")
    return overlaps


def main() -> int:
    ap = argparse.ArgumentParser(
        description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter
    )
    ap.add_argument(
        "--slide", type=int, default=None, help="Only dump this slide (1-based)."
    )
    ap.add_argument(
        "--overlap", action="store_true", help="Flag overlapping shape pairs per slide."
    )
    ap.add_argument(
        "--deck", type=Path, default=DECK, help=f"Path to .pptx (default: {DECK})."
    )
    args = ap.parse_args()

    if not args.deck.exists():
        print(f"✗ deck not found: {args.deck}", file=sys.stderr)
        return 1

    prs = Presentation(str(args.deck))
    total = len(prs.slides)
    print(f"BeThere shape inventory — {args.deck.name}")
    print(f'Canvas: {CANVAS_W_IN}" × {CANVAS_H_IN}"  ·  {total} slides')

    total_overlaps = 0
    for idx, slide in enumerate(prs.slides, start=1):
        if args.slide is not None and idx != args.slide:
            continue
        overlaps = dump_slide(slide, idx, flag_overlap=args.overlap)
        total_overlaps += len(overlaps)

    if args.overlap:
        print(f"\n{'=' * 78}")
        print(f"TOTAL OVERLAPS across dumped slides: {total_overlaps}")
        print(f"(Note: text_box vs text_box overlaps are common and usually")
        print(f" benign — the textbox spans wider than the rendered text.)")

    return 0


if __name__ == "__main__":
    sys.exit(main())
